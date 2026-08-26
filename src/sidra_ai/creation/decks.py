"""Build a deck whose every number came from somewhere.

"デッキを作って" is the request where a language model is most dangerous.
A slide wants a figure in it, the shape of the sentence pulls one out, and
the deck reads as authoritative precisely because it is formatted like one.
So this generator has a rule that overrides everything else:

    **A number appears on a slide only if it was retrieved. Otherwise the
    slide shows 〔社長が埋める欄〕 and says which slide is unfilled.**

That is not a limitation to apologise for - it is the feature. A deck with
three blanks the owner fills in is usable. A deck with three invented figures
is worse than no deck, and looks identical from the outside.

Everything else follows the game generator next door: the deck is a template
first, correct with no model and no network, and a local model may only
overlay wording through :meth:`GeneratedDeck.with_copy`. The HTML file is the
artifact that always works; ``python-pptx`` produces a .pptx as well when it
happens to be importable, and is never added as a dependency.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass, field, replace
from datetime import datetime, timezone
from html import escape
from pathlib import Path

#: Shared with the game generator so a deck and a game made by the same tool
#: look like they came from the same place.
from sidra_ai.creation.evidence import NUMBER, Fact
from sidra_ai.creation.themes import Theme, select_theme

#: What an unfilled slot says. Kept as one constant because both the renderer
#: and the validator have to agree on it exactly: the validator counts these
#: to report how much of the deck is still the owner's to write.
BLANK = "〔社長が埋める欄〕"

#: The evidence type and the number pattern live in ``creation.evidence`` now
#: that the router carries facts to any generator. Re-exported here because a
#: deck is still the main thing built from them.
_NUMBER = NUMBER


@dataclass(frozen=True)
class Slide:
    title: str
    bullets: tuple[str, ...] = ()
    sources: tuple[str, ...] = ()

    @property
    def blanks(self) -> int:
        return sum(1 for bullet in self.bullets if BLANK in bullet)


@dataclass(frozen=True)
class DeckOutline:
    """The sections a deck has, before any content is available.

    Fixed rather than chosen by a model: the shape of a pitch is not the
    interesting part, and a model that invents section names also invents
    what belongs under them.
    """

    key: str
    default_title: str
    sections: tuple[str, ...]


OUTLINES: dict[str, DeckOutline] = {
    "pitch": DeckOutline(
        key="pitch",
        default_title="SIDRA AI のご提案",
        sections=("課題", "解決", "根拠となる数字", "次の一歩"),
    ),
    "status": DeckOutline(
        key="status",
        default_title="進捗報告",
        sections=("いま出来ること", "測った数字", "残っていること", "判断が要る点"),
    ),
}


@dataclass(frozen=True)
class GeneratedDeck:
    outline: str
    title: str
    slides: tuple[Slide, ...]
    html: str
    #: Slide titles that still contain a blank, in order. Carried on the
    #: object rather than recomputed, so the summary an operator reads and
    #: the artifact they open cannot disagree.
    unfilled: tuple[str, ...] = ()

    def with_copy(self, *, title: str = "") -> "GeneratedDeck":
        """Overlay model-written wording on a deck that is already correct.

        Only the title. A model may not touch bullets, because a bullet is
        where a number would be reworded into existence.
        """

        if not title.strip():
            return self
        clean = " ".join(title.split())[:80]
        html = self.html.replace(
            f"<title>{escape(self.title)}</title>", f"<title>{escape(clean)}</title>"
        ).replace(f"<h1>{escape(self.title)}</h1>", f"<h1>{escape(clean)}</h1>")
        return replace(self, title=clean, html=html)


def choose_outline(request: str) -> str:
    """Pick a shape from the request, defaulting to a pitch."""

    text = request.casefold()
    for word in ("進捗", "報告", "status", "レポート", "週報"):
        if word.casefold() in text:
            return "status"
    return "pitch"


def _title_from(request: str, fallback: str) -> str:
    """Use the operator's own words for the title when there are any.

    Cut at the making-verb, because "営業用のデッキを作って" titles better as
    "営業用のデッキ" than as the whole imperative.
    """

    stripped = re.split(r"を?(?:作って|作成して|生成して|つくって)", request)[0].strip()
    stripped = " ".join(stripped.split())
    return stripped[:60] or fallback


#: What a passage has to contain to belong under a section. Literal cues, no
#: model: a matcher clever enough to place any passage under any heading is a
#: matcher that fills every slide, and a slide filled with a sentence that
#: does not answer it reads as an answer.
#:
#: Measured, not guessed: matching on the heading word alone ("課題" inside the
#: passage) retrieved five passages from this repository and placed none of
#: them, because retrieved prose almost never repeats the heading it belongs
#: under. Cues are the smallest fix that keeps the conservative direction -
#: a passage with no cue is still left out.
SECTION_CUES: dict[str, tuple[str, ...]] = {
    "課題": ("課題", "問題", "できない", "困", "ギャップ", "未対応", "gap", "problem"),
    "解決": ("解決", "対応", "実装", "できる", "提供", "returns", "provides", "solution"),
    "根拠となる数字": (),  # decided by the presence of a number, see below
    "次の一歩": ("次", "予定", "残り", "todo", "roadmap", "backlog", "next"),
    "いま出来ること": ("できる", "対応", "提供", "supports", "provides"),
    "測った数字": (),  # same rule as 根拠となる数字
    "残っていること": ("残", "未", "todo", "backlog", "gap"),
    "判断が要る点": ("判断", "要判断", "決め", "decision", "trade-off"),
}

#: Sections whose whole point is a figure. A passage belongs here when it
#: carries one - which also means every number on the slide came from the
#: evidence by construction, the property ``validate_deck`` re-checks.
_NUMERIC_SECTIONS = ("根拠となる数字", "測った数字")


def _matches(section: str, fact: Fact) -> bool:
    if section in _NUMERIC_SECTIONS:
        return fact.mentions_number()
    if section in fact.text:
        return True
    text = fact.text.casefold()
    return any(cue.casefold() in text for cue in SECTION_CUES.get(section, ()))


def _bullets_for(section: str, facts: list[Fact]) -> tuple[tuple[str, ...], tuple[str, ...]]:
    """Fill one section from the facts that belong under it, or leave it blank."""

    hits = [fact for fact in facts if _matches(section, fact)]
    if not hits:
        return (f"{BLANK}",), ()
    bullets = tuple(" ".join(fact.text.split())[:120] for fact in hits[:3])
    sources = tuple(dict.fromkeys(fact.source for fact in hits[:3]))
    return bullets, sources


def build_slides(outline: DeckOutline, facts: list[Fact]) -> tuple[Slide, ...]:
    """One slide per section. Sections with no evidence keep their blank."""

    slides: list[Slide] = []
    for section in outline.sections:
        bullets, sources = _bullets_for(section, facts)
        slides.append(Slide(title=section, bullets=bullets, sources=sources))
    return tuple(slides)


def _no_external_assets(html: str) -> bool:
    for match in re.finditer(r"""(?:src|href)\s*=\s*["']([^"']+)["']""", html):
        if match.group(1).strip().lower().startswith(("http://", "https://", "//")):
            return False
    return "@import" not in html


def _render(title: str, slides: tuple[Slide, ...], theme: Theme) -> str:
    t = theme.tokens
    blocks = []
    for index, slide in enumerate(slides, start=1):
        bullets = "".join(f"<li>{escape(b)}</li>" for b in slide.bullets)
        sources = (
            "<p class='src'>出典: " + escape(" / ".join(slide.sources)) + "</p>"
            if slide.sources
            else "<p class='src blank'>出典なし - この欄は埋まっていません</p>"
        )
        blocks.append(
            f"<section class='slide'><p class='no'>{index}/{len(slides)}</p>"
            f"<h2>{escape(slide.title)}</h2><ul>{bullets}</ul>{sources}</section>"
        )
    return f"""<!doctype html>
<html lang="ja"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escape(title)}</title>
<style>
:root{{color-scheme:{t["scheme"]}}}
body{{margin:0;background:{t["bg"]};color:{t["text"]};
 font-family:system-ui,"Hiragino Kaku Gothic ProN","Noto Sans JP",sans-serif}}
main{{max-width:900px;margin:0 auto;padding:32px 20px 56px}}
h1{{font-size:24px;margin:0 0 24px;letter-spacing:.01em}}
.slide{{background:{t["surface"]};border:1px solid {t["border"]};
 border-radius:{t["radius"]};padding:22px 24px;margin:0 0 18px}}
.slide .no{{margin:0 0 6px;font-size:12px;color:{t["muted"]}}}
.slide h2{{margin:0 0 12px;font-size:19px;color:{t["accent"]}}}
.slide ul{{margin:0;padding-left:20px;line-height:1.75}}
.slide .src{{margin:14px 0 0;font-size:12px;color:{t["muted"]}}}
.slide .blank{{color:{t["alert"]}}}
footer{{margin-top:24px;border-top:1px solid {t["border"]};padding-top:14px;
 font-size:12px;color:{t["muted"]}}}
</style></head>
<body><main>
<h1>{escape(title)}</h1>
{"".join(blocks)}
<footer>SIDRA AI が生成。数字は索引した文書から引いたものだけを載せ、
根拠が無い欄は {escape(BLANK)} のまま残しています（推測で埋めません）。</footer>
</main></body></html>
"""


def generate_deck(
    request: str,
    *,
    facts: list[Fact] | None = None,
    outline: str = "",
) -> GeneratedDeck:
    """Build a deck from the request and whatever evidence was retrieved.

    Never raises on wording, and never invents content: with no facts at all
    the deck still renders, entirely in blanks, which is an honest artifact.
    """

    key = outline or choose_outline(request)
    if key not in OUTLINES:
        raise KeyError(f"unknown deck outline: {key!r}")
    spec = OUTLINES[key]
    slides = build_slides(spec, list(facts or []))
    title = _title_from(request, spec.default_title)
    html = _render(title, slides, select_theme(request))
    unfilled = tuple(slide.title for slide in slides if slide.blanks)
    return GeneratedDeck(key, title, slides, html, unfilled)


def save_deck(deck: GeneratedDeck, data_dir: str | Path, *, now: datetime | None = None) -> Path:
    """Write the artifact locally. Nothing leaves the machine."""

    stamp = (now or datetime.now(timezone.utc)).strftime("%Y%m%dT%H%M%SZ")
    directory = Path(data_dir) / "artifacts"
    directory.mkdir(parents=True, exist_ok=True)
    path = directory / f"deck-{deck.outline}-{stamp}.html"
    path.write_text(deck.html, encoding="utf-8")
    return path


def save_pptx(deck: GeneratedDeck, path: str | Path) -> tuple[bool, str]:
    """Also write a .pptx when ``python-pptx`` is importable.

    Optional by design. The HTML file is the artifact that always exists;
    a slide format that needs a package the operator may not have cannot be
    the thing that decides whether a deck was produced. Returns whether it
    was written and why not, so a caller reports the truth rather than
    silently producing one format while claiming another.
    """

    try:
        from pptx import Presentation  # type: ignore[import-not-found]
        from pptx.util import Pt  # type: ignore[import-not-found]
    except Exception as exc:  # noqa: BLE001 - absence is the expected case
        return False, f"python-pptx not available: {type(exc).__name__}"

    presentation = Presentation()
    title_layout, bullet_layout = presentation.slide_layouts[0], presentation.slide_layouts[1]
    opening = presentation.slides.add_slide(title_layout)
    opening.shapes.title.text = deck.title

    for slide in deck.slides:
        rendered = presentation.slides.add_slide(bullet_layout)
        rendered.shapes.title.text = slide.title
        frame = rendered.placeholders[1].text_frame
        frame.text = slide.bullets[0] if slide.bullets else BLANK
        for bullet in slide.bullets[1:]:
            paragraph = frame.add_paragraph()
            paragraph.text = bullet
            paragraph.font.size = Pt(16)

    Path(path).parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))
    return True, "python-pptx"


# -------------------------------------------------------------- validation


def validate_deck(deck: GeneratedDeck, facts: list[Fact] | None = None) -> dict:
    """Report every reason the deck is not usable, not just the first.

    The load-bearing check is the last one. A deck that parses, has slides
    and fetches nothing is still a failure if a figure on it appears in no
    retrieved fact, because that figure came from somewhere other than the
    corpus - which, in a generator, means it was made up.
    """

    from html.parser import HTMLParser

    failures: list[str] = []

    parser = HTMLParser(convert_charrefs=True)
    try:
        parser.feed(deck.html)
        parser.close()
    except Exception as exc:  # noqa: BLE001 - a parse failure is a finding
        failures.append(f"html: {type(exc).__name__}: {exc}")

    if not deck.slides:
        failures.append("no slides")
    if deck.html.count("<section class='slide'>") != len(deck.slides):
        failures.append("rendered slide count does not match the outline")
    if not _no_external_assets(deck.html):
        failures.append("references an external asset")

    evidence = " ".join(fact.text for fact in (facts or []))
    unsourced: list[str] = []
    for slide in deck.slides:
        for bullet in slide.bullets:
            if BLANK in bullet:
                continue
            for number in _NUMBER.findall(bullet):
                token = number.strip()
                if token and token not in evidence:
                    unsourced.append(f"{slide.title}: {token}")
    if unsourced:
        failures.append("numbers not present in the evidence: " + ", ".join(unsourced[:5]))

    return {
        "usable": not failures,
        "failures": failures,
        "slides": len(deck.slides),
        "unfilled": list(deck.unfilled),
    }


def report(deck: GeneratedDeck) -> str:
    return json.dumps(
        {
            "outline": deck.outline,
            "title": deck.title,
            "slides": len(deck.slides),
            "unfilled": list(deck.unfilled),
        },
        ensure_ascii=False,
    )
